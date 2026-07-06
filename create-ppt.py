from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
DARK_BG = RGBColor(15, 23, 42)
BLUE = RGBColor(37, 99, 235)
AMBER = RGBColor(245, 158, 11)
GREEN = RGBColor(5, 150, 105)
PURPLE = RGBColor(124, 58, 237)
WHITE = RGBColor(255, 255, 255)
GRAY = RGBColor(148, 163, 184)
PINK = RGBColor(236, 72, 153)
DARK_CARD = RGBColor(30, 41, 59)
BORDER = RGBColor(51, 65, 85)
DARK_ROW1 = RGBColor(20, 30, 50)
DARK_ROW2 = RGBColor(25, 35, 55)

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_text(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Calibri'
        p.space_after = Pt(8)
    return txBox

def add_glow(slide, left, top, w, h, color, alpha_val='12000'):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    spPr = shape._element.find(qn('p:spPr'))
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is None:
        solidFill = spPr.find('.//' + qn('a:solidFill'))
    if solidFill is not None:
        srgb = solidFill.find(qn('a:srgbClr'))
        if srgb is not None:
            a = etree.SubElement(srgb, qn('a:alpha'))
            a.set('val', alpha_val)
    shape.line.fill.background()
    return shape

# ==================== SLIDE 1: TITLE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_glow(slide, Inches(9), Inches(-2), Inches(6), Inches(6), BLUE, '15000')
add_glow(slide, Inches(-2), Inches(4), Inches(5), Inches(5), PURPLE, '12000')
add_text(slide, Inches(1), Inches(1.5), Inches(8), Inches(1), 'BILL BY BILLU', 20, True, AMBER)
add_text(slide, Inches(1), Inches(2.2), Inches(9), Inches(1.5), 'AI-Powered GST Invoicing\nfor Indian Businesses', 48, True, WHITE)
add_text(slide, Inches(1), Inches(4.2), Inches(8), Inches(0.8), 'Client Proposal & Product Overview', 22, False, GRAY)
add_text(slide, Inches(1), Inches(5.5), Inches(6), Inches(0.5), 'www.billbybillu.in  •  amaradworld@gmail.com  •  Play Store: com.BillbBillu', 14, False, GRAY)

# ==================== SLIDE 2: PROBLEM ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'THE PROBLEM', 14, True, AMBER)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'Indian SMBs Struggle with GST Compliance', 40, True, WHITE)

problems = [
    ('Manual Invoicing', 'Spreadsheets, paper bills, WhatsApp messages — no structure'),
    ('GST Filing Hassle', 'GSTR-1/GSTR-3B reports take hours to compile manually'),
    ('Payment Tracking', 'No visibility into who paid, who owes — cash flow gaps'),
    ('Language Barrier', 'Most tools are English-only — team can\'t use them'),
    ('No Mobile Access', 'Desktop-only software — can\'t create invoices on the go'),
]

for i, (title, desc) in enumerate(problems):
    x = Inches(0.8) if i < 3 else Inches(6.8)
    y = Inches(2.5) + Inches(1.3) * (i % 3)
    add_shape(slide, x, y, Inches(5.5), Inches(1.1), DARK_CARD, BORDER)
    add_text(slide, x + Inches(0.3), y + Inches(0.15), Inches(4.8), Inches(0.4), f'❌  {title}', 18, True, WHITE)
    add_text(slide, x + Inches(0.3), y + Inches(0.55), Inches(4.8), Inches(0.4), desc, 13, False, GRAY)

# ==================== SLIDE 3: SOLUTION ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'OUR SOLUTION', 14, True, GREEN)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'Bill By Billu — One App, Everything You Need', 40, True, WHITE)

solutions = [
    ('🤖', 'AI Invoice Creator', 'Voice, photo, or chat — AI creates\nGST invoices in seconds', BLUE),
    ('📊', 'GST Reports', 'Auto-generated GSTR-1 & GSTR-3B,\nready to file on gst.gov.in', GREEN),
    ('💬', 'WhatsApp Share', 'Send invoices to customers\nin one tap with payment links', RGBColor(37, 211, 102)),
    ('💳', 'UPI Payments', 'Auto QR code with exact\namount on every invoice', PURPLE),
    ('📄', '3 Templates', 'Classic, Modern & Compact\ndesigns with your logo', AMBER),
    ('🌐', '9 Languages', 'Invoice in your\ncustomer\'s language', PINK),
]

for i, (icon, title, desc, accent) in enumerate(solutions):
    col = i % 3
    row = i // 3
    x = Inches(0.8) + Inches(4.1) * col
    y = Inches(2.5) + Inches(2.3) * row
    add_shape(slide, x, y, Inches(3.8), Inches(2.0), DARK_CARD, BORDER)
    add_circle(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.7), accent)
    add_text(slide, x + Inches(0.3), y + Inches(0.35), Inches(0.7), Inches(0.6), icon, 24, False, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(1.2), y + Inches(0.35), Inches(2.4), Inches(0.4), title, 18, True, WHITE)
    add_text(slide, x + Inches(0.3), y + Inches(1.2), Inches(3.2), Inches(0.7), desc, 13, False, GRAY)

# ==================== SLIDE 4: FEATURES ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'KEY FEATURES', 14, True, BLUE)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'Built for How Indian Businesses Actually Work', 40, True, WHITE)

features = [
    'AI-powered invoice creation via voice, photo OCR, or chat',
    'GSTR-1 & GSTR-3B auto-generated — export CSV/JSON',
    'WhatsApp sharing with payment link on every invoice',
    'UPI QR code auto-generated with exact invoice amount',
    '3 premium invoice templates (Classic, Modern, Compact)',
    'Custom logo & bank details on every invoice',
    'Credit notes, debit notes & recurring invoices',
    'Multi-currency support with auto-conversion',
    'Customer & product database with search',
    'Expense tracking with AI categorization (12 categories)',
    'Real-time dashboard with revenue analytics',
    'Role-based access: Owner, Accountant, Viewer',
]
add_bullet_list(slide, Inches(0.8), Inches(2.3), Inches(5.8), Inches(4.5), [f'✓  {f}' for f in features[:6]], 15, WHITE)
add_bullet_list(slide, Inches(6.8), Inches(2.3), Inches(5.8), Inches(4.5), [f'✓  {f}' for f in features[6:]], 15, WHITE)

# ==================== SLIDE 5: AI CREATOR ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'AI POWERED', 14, True, AMBER)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'AI Invoice Creator — Just Describe, We Create', 40, True, WHITE)

add_text(slide, Inches(0.8), Inches(2.5), Inches(5), Inches(0.5), 'How It Works', 22, True, AMBER)
steps = [
    ('1', 'Speak or Type', '"Sold 10 t-shirts at ₹450 to Rajesh, GST 5%"'),
    ('2', 'AI Parses Everything', 'Items, quantities, rates, GST — all extracted automatically'),
    ('3', 'Invoice Created', 'Professional GST invoice with all details in seconds'),
    ('4', 'Share Instantly', 'Send via WhatsApp, email, or download PDF'),
]
for i, (num, title, desc) in enumerate(steps):
    y = Inches(3.2) + Inches(0.9) * i
    add_circle(slide, Inches(0.8), y, Inches(0.5), BLUE)
    add_text(slide, Inches(0.8), y + Inches(0.05), Inches(0.5), Inches(0.4), num, 16, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, Inches(1.5), y, Inches(4), Inches(0.3), title, 16, True, WHITE)
    add_text(slide, Inches(1.5), y + Inches(0.3), Inches(4.5), Inches(0.4), desc, 12, False, GRAY)

add_text(slide, Inches(7), Inches(2.5), Inches(5), Inches(0.5), '3 Ways to Create', 22, True, AMBER)
methods = [
    ('🎤', 'Voice Input', 'Speak naturally in Hindi or English', BLUE),
    ('📸', 'Photo OCR', 'Snap a photo of a handwritten bill', GREEN),
    ('💬', 'Chat Input', 'Type like you\'re messaging a friend', PURPLE),
]
for i, (icon, title, desc, color) in enumerate(methods):
    y = Inches(3.2) + Inches(1.1) * i
    add_shape(slide, Inches(7), y, Inches(5.5), Inches(0.9), DARK_CARD, BORDER)
    add_text(slide, Inches(7.3), y + Inches(0.1), Inches(0.5), Inches(0.5), icon, 24, False, WHITE)
    add_text(slide, Inches(8), y + Inches(0.1), Inches(4), Inches(0.3), title, 16, True, WHITE)
    add_text(slide, Inches(8), y + Inches(0.45), Inches(4), Inches(0.3), desc, 12, False, GRAY)

# ==================== SLIDE 6: GST ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'GST COMPLIANCE', 14, True, GREEN)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'GSTR-1 & GSTR-3B — Auto-Generated, Ready to File', 40, True, WHITE)

add_shape(slide, Inches(0.8), Inches(2.5), Inches(5.8), Inches(4.2), DARK_CARD, BORDER)
add_text(slide, Inches(1.2), Inches(2.7), Inches(5), Inches(0.5), '📊  GSTR-1 (Outward Supplies)', 20, True, GREEN)
add_bullet_list(slide, Inches(1.2), Inches(3.4), Inches(5), Inches(3),
    ['→  B2B invoice details with HSN codes', '→  Auto-calculated CGST, SGST, IGST',
     '→  Intra-state vs inter-state detection', '→  HSN-wise summary of goods/services',
     '→  Export in JSON format for gst.gov.in', '→  CSV export for Excel/spreadsheet use'], 13, WHITE)

add_shape(slide, Inches(7), Inches(2.5), Inches(5.8), Inches(4.2), DARK_CARD, BORDER)
add_text(slide, Inches(7.4), Inches(2.7), Inches(5), Inches(0.5), '📋  GSTR-3B (Monthly Return)', 20, True, GREEN)
add_bullet_list(slide, Inches(7.4), Inches(3.4), Inches(5), Inches(3),
    ['→  Auto-populated from GSTR-1 data', '→  Tax liability calculation',
     '→  Input tax credit summary', '→  Net tax payable computation',
     '→  Period-wise reporting (monthly)', '→  One-click export for filing'], 13, WHITE)

# ==================== SLIDE 7: PRICING ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'PRICING', 14, True, PURPLE)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'Simple, Transparent Pricing', 40, True, WHITE)

plans = [
    ('FREE', '₹0', 'forever', ['5 invoices/month', '3 invoice templates', 'PDF export', 'Customer database', 'Basic dashboard'], GRAY),
    ('STARTER', '₹199', '/month', ['100 invoices/month', 'Custom logo on invoices', 'Credit/Debit notes', 'Recurring invoices', 'UPI payment links', 'Priority support'], BLUE),
    ('GROWTH', '₹399', '/month', ['500 invoices/month', 'AI Invoice Creator', 'GSTR-1 & GSTR-3B', 'WhatsApp sharing', 'Expense tracking', 'Multi-user access', 'API access'], AMBER),
    ('ENTERPRISE', 'Custom', 'contact us', ['Unlimited invoices', 'White-label option', 'Dedicated support', 'Custom integrations', 'SLA guarantee'], PURPLE),
]

for i, (name, price, period, features, accent) in enumerate(plans):
    x = Inches(0.5) + Inches(3.15) * i
    y = Inches(2.5)
    w = Inches(3.0)
    h = Inches(4.5)
    
    if name == 'GROWTH':
        add_shape(slide, x - Inches(0.05), y - Inches(0.05), w + Inches(0.1), h + Inches(0.1), accent)
        badge = add_shape(slide, x + Inches(0.7), y - Inches(0.25), Inches(1.6), Inches(0.35), AMBER)
        add_text(slide, x + Inches(0.7), y - Inches(0.23), Inches(1.6), Inches(0.3), 'MOST POPULAR', 10, True, DARK_BG, PP_ALIGN.CENTER)
    
    add_shape(slide, x, y, w, h, DARK_CARD, BORDER)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.4), Inches(0.4), name, 16, True, accent, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), y + Inches(0.8), Inches(2.4), Inches(0.6), price, 36, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), y + Inches(1.3), Inches(2.4), Inches(0.3), period, 12, False, GRAY, PP_ALIGN.CENTER)
    
    add_rect(slide, x + Inches(0.3), y + Inches(1.7), Inches(2.4), Pt(1), BORDER)
    add_bullet_list(slide, x + Inches(0.3), y + Inches(1.9), Inches(2.4), Inches(2.5), [f'✓  {f}' for f in features], 11, WHITE)

add_text(slide, Inches(0.8), Inches(7.1), Inches(11), Inches(0.4), '💰  Annual plans save 21% — Starter ₹1,887/yr, Growth ₹3,783/yr', 14, False, AMBER)

# ==================== SLIDE 8: TECH ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'TECHNOLOGY', 14, True, BLUE)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'Modern Tech Stack, Built to Scale', 40, True, WHITE)

tech_items = [
    ('Frontend', 'React + Vite, TailwindCSS, Capacitor (Android)', BLUE),
    ('Backend', 'Node.js + Express, Prisma ORM', GREEN),
    ('Database', 'PostgreSQL (managed)', PURPLE),
    ('AI Engine', 'Custom NLP + Google Gemini API', AMBER),
    ('PDF Generation', 'PDFKit with custom templates', RGBColor(37, 211, 102)),
    ('Payments', 'Razorpay (UPI, Cards, Netbanking)', PINK),
    ('Auth', 'JWT + Google OAuth 2.0', BLUE),
    ('Hosting', 'Render (backend) + Vercel (frontend)', GREEN),
    ('Mobile', 'Capacitor → Google Play Store', PURPLE),
    ('Email', 'SMTP (Zoho Mail / Titan)', AMBER),
    ('GST Integration', 'Auto HSN code detection, GSTIN validation', RGBColor(37, 211, 102)),
    ('Security', 'HMAC webhooks, rate limiting, RBAC', PINK),
]

for i, (label, desc, color) in enumerate(tech_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.2) * col
    y = Inches(2.3) + Inches(0.75) * row
    add_rect(slide, x, y + Inches(0.05), Inches(0.12), Inches(0.45), color)
    add_text(slide, x + Inches(0.3), y, Inches(2), Inches(0.3), label, 14, True, color)
    add_text(slide, x + Inches(0.3), y + Inches(0.3), Inches(5.5), Inches(0.3), desc, 12, False, GRAY)

# ==================== SLIDE 9: SECURITY ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'SECURITY', 14, True, GREEN)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'Enterprise-Grade Security', 40, True, WHITE)

security_items = [
    ('🔒', 'JWT Authentication', 'Secure token-based auth with HS256 pinning', BLUE),
    ('🛡️', 'Rate Limiting', '500 requests/15min per IP, DDoS protection', GREEN),
    ('🔐', 'HMAC Webhooks', 'Razorpay webhook verification with anti-replay', PURPLE),
    ('👤', 'RBAC', 'Role-based access: Owner, Accountant, Viewer', AMBER),
    ('🌐', 'CORS Protection', 'Strict origin whitelisting for production', RGBColor(37, 211, 102)),
    ('🔑', 'Secret Management', 'Environment variables, no hardcoded secrets', PINK),
    ('📱', 'Capacitor Security', 'Secure token storage via native API', BLUE),
    ('🛡️', 'Input Validation', 'Zod schema validation on all endpoints', GREEN),
]

for i, (icon, title, desc, color) in enumerate(security_items):
    col = i % 2
    row = i // 2
    x = Inches(0.8) + Inches(6.2) * col
    y = Inches(2.3) + Inches(1.15) * row
    add_shape(slide, x, y, Inches(5.8), Inches(0.95), DARK_CARD, BORDER)
    add_text(slide, x + Inches(0.3), y + Inches(0.1), Inches(0.5), Inches(0.5), icon, 22, False, WHITE)
    add_text(slide, x + Inches(1), y + Inches(0.1), Inches(4.5), Inches(0.3), title, 16, True, WHITE)
    add_text(slide, x + Inches(1), y + Inches(0.45), Inches(4.5), Inches(0.35), desc, 12, False, GRAY)

# ==================== SLIDE 10: COMPETITIVE ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'WHY US', 14, True, AMBER)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'What Makes Bill By Billu Different', 40, True, WHITE)

headers = ['Feature', 'Bill By Billu', 'Zoho Invoice', 'Vyapar', 'Khatabook']
col_widths = [Inches(2.8), Inches(2.0), Inches(2.0), Inches(2.0), Inches(2.0)]
col_x = [Inches(0.8)]
for w in col_widths[:-1]:
    col_x.append(col_x[-1] + w + Inches(0.1))

for j, header in enumerate(headers):
    bg = GREEN if j == 1 else DARK_CARD
    add_shape(slide, col_x[j], Inches(2.3), col_widths[j], Inches(0.5), bg)
    add_text(slide, col_x[j], Inches(2.32), col_widths[j], Inches(0.45), header, 13, True, WHITE, PP_ALIGN.CENTER)

rows_data = [
    ['AI Invoice Creation', '✅', '❌', '❌', '❌'],
    ['WhatsApp Sharing', '✅', '✅', '✅', '✅'],
    ['UPI QR on Invoice', '✅', '❌', '✅', '❌'],
    ['GST Reports (GSTR-1/3B)', '✅', '✅', '✅', '❌'],
    ['9 Indian Languages', '✅', '❌', '✅', '✅'],
    ['Free Plan', '5 invoices', '1000 inv', 'Limited', 'Limited'],
    ['Mobile App', '✅ Android', '✅', '✅', '✅'],
    ['Starting Price', '₹199/mo', '₹749/mo', '₹499/mo', '₹399/mo'],
]

for i, row in enumerate(rows_data):
    y = Inches(2.9) + Inches(0.5) * i
    bg = DARK_ROW1 if i % 2 == 0 else DARK_ROW2
    for j, cell in enumerate(row):
        add_shape(slide, col_x[j], y, col_widths[j], Inches(0.45), bg, BORDER if j == 1 else None)
        is_bold = j == 0 or (j == 1 and cell == '✅')
        cell_color = GREEN if j == 1 and cell == '✅' else (WHITE if j == 0 else GRAY)
        add_text(slide, col_x[j], y + Inches(0.05), col_widths[j], Inches(0.35), cell, 12, is_bold, cell_color, PP_ALIGN.CENTER)

# ==================== SLIDE 11: ROADMAP ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_text(slide, Inches(0.8), Inches(0.5), Inches(5), Inches(0.6), 'ROADMAP', 14, True, BLUE)
add_text(slide, Inches(0.8), Inches(1.1), Inches(10), Inches(1), 'What\'s Coming Next', 40, True, WHITE)

phases = [
    ('Q3 2026', 'LAUNCHED', ['AI Invoice Creator', 'GSTR-1 & GSTR-3B reports', 'WhatsApp sharing', '3 premium templates', 'UPI QR on invoices', 'Android Play Store'], GREEN),
    ('Q4 2026', 'IN PROGRESS', ['iOS app (Capacitor)', 'Multi-business support', 'Inventory management', 'E-invoicing (IRN)', 'Tally/Busy integration', 'Advanced analytics'], AMBER),
    ('Q1 2027', 'PLANNED', ['E-way bill generation', 'GSTN API integration', 'Multi-user permissions', 'API for third-party apps', 'White-label solution', 'TDS/TCS support'], PURPLE),
]

for i, (quarter, status, items, color) in enumerate(phases):
    x = Inches(0.8) + Inches(4.1) * i
    add_shape(slide, x, Inches(2.3), Inches(3.8), Inches(4.8), DARK_CARD, BORDER)
    add_shape(slide, x + Inches(0.3), Inches(2.5), Inches(1.6), Inches(0.35), color)
    add_text(slide, x + Inches(0.3), Inches(2.52), Inches(1.6), Inches(0.3), status, 10, True, WHITE, PP_ALIGN.CENTER)
    add_text(slide, x + Inches(0.3), Inches(3.0), Inches(3.2), Inches(0.4), quarter, 22, True, WHITE)
    add_bullet_list(slide, x + Inches(0.3), Inches(3.5), Inches(3.2), Inches(3.3), [f'→  {item}' for item in items], 12, WHITE)

# ==================== SLIDE 12: CTA ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(slide, DARK_BG)
add_glow(slide, Inches(4), Inches(1), Inches(6), Inches(6), BLUE, '10000')

add_text(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1), 'Let\'s Build Together', 56, True, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(2), Inches(3), Inches(9), Inches(0.8), 'Bill By Billu is ready to transform how Indian businesses handle invoicing & GST compliance.', 20, False, GRAY, PP_ALIGN.CENTER)

add_shape(slide, Inches(3.5), Inches(4.2), Inches(6), Inches(2.5), DARK_CARD, BORDER)
add_text(slide, Inches(3.5), Inches(4.4), Inches(6), Inches(0.5), 'Get in Touch', 24, True, AMBER, PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(5.0), Inches(6), Inches(0.4), '🌐  www.billbybillu.in', 16, False, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(5.4), Inches(6), Inches(0.4), '📧  amaradworld@gmail.com', 16, False, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(5.8), Inches(6), Inches(0.4), '📱  Play Store: Bill By Billu', 16, False, WHITE, PP_ALIGN.CENTER)
add_text(slide, Inches(3.5), Inches(6.2), Inches(6), Inches(0.4), '💼  github.com/amaradworld/bill-by-billu', 14, False, GRAY, PP_ALIGN.CENTER)

# ==================== SAVE ====================
output_path = r'C:\Users\alokg\bill-by-billu\BillByBillu-ClientProposal.pptx'
prs.save(output_path)
print(f'Saved: {output_path}')
print(f'Slides: {len(prs.slides)}')
